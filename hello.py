# flight_difficulty.py
import pandas as pd
import numpy as np
from pathlib import Path
YOUR_NAME = "chaitanyakumar"
INPUT_DIR = Path(".")
FLIGHTS_CSV = "./data/Flight%20Level%20Data.csv"
PNR_CSV = "./data/PNR+Flight+Level+Data.csv"
PNR_REMARKS_CSV = "./data/PNR%20Remark%20Level%20Data.csv"
BAGS_CSV = "./data/Bag+Level+Data.csv"
AIRPORTS_CSV = "./data/Airports%20Data.csv"
OUTPUT_CSV = f"./data/test_{YOUR_NAME}.csv"
import sys
import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
def safe_div(a, b):
    return a / b.replace(0, np.nan) if hasattr(b, "replace") else np.divide(a, b, out=np.zeros_like(a, dtype=float), where=b!=0)
def check_input_files(paths):
    missing = [p for p in paths if not Path(p).exists()]
    if missing:
        print("The following required input files are missing:")
        for m in missing:
            print(f"  - {m}")
        print("\nPlease ensure the files exist at the paths above (relative to this script) or update the path constants.")
        raise SystemExit(2)
def ensure_columns(df, cols_with_defaults):
    for col, default in cols_with_defaults.items():
        if col not in df.columns:
            df[col] = default
def main():
    check_input_files([FLIGHTS_CSV, PNR_CSV, PNR_REMARKS_CSV, BAGS_CSV, AIRPORTS_CSV])
    flights = pd.read_csv(FLIGHTS_CSV, parse_dates=["scheduled_departure_datetime_local", "scheduled_arrival_datetime_local", "actual_departure_datetime_local", "actual_arrival_datetime_local"])
    pnr = pd.read_csv(PNR_CSV, parse_dates=["pnr_creation_date"])
    pnr_remarks = pd.read_csv(PNR_REMARKS_CSV, parse_dates=["pnr_creation_date"])
    bags = pd.read_csv(BAGS_CSV, parse_dates=["bag_tag_issue_date"])
    airports = pd.read_csv(AIRPORTS_CSV)
    ensure_columns(pnr, {
        "total_pax": 0,
        "lap_child_count": 0,
        "basic_economy_pax": 0,
        "is_stroller_user": 0,
        "company_id": np.nan,
        "flight_number": np.nan,
        "scheduled_departure_date_local": pd.NaT
    })
    ensure_columns(pnr_remarks, {
        "special_service_request": "",
        "company_id": np.nan,
        "flight_number": np.nan,
        "scheduled_departure_date_local": pd.NaT
    })
    ensure_columns(bags, {
        "bag_type": "",
        "company_id": np.nan,
        "flight_number": np.nan,
        "scheduled_departure_date_local": pd.NaT
    })
    for df in [pnr, pnr_remarks, bags]:
        if "scheduled_departure_date_local" not in df.columns or df["scheduled_departure_date_local"].isnull().all():
            for col in df.columns:
                if "scheduled" in col and "datetime" in col:
                    try:
                        df["scheduled_departure_date_local"] = pd.to_datetime(df[col]).dt.date
                        break
                    except Exception:
                        continue
    for df in [flights, pnr, pnr_remarks, bags]:
        if "scheduled_departure_date_local" in df.columns:
            df["scheduled_departure_date_local"] = pd.to_datetime(df["scheduled_departure_date_local"]).dt.date
    flights["departure_delay_minutes"] = (flights["actual_departure_datetime_local"] - flights["scheduled_departure_datetime_local"]).dt.total_seconds() / 60
    avg_delay = flights["departure_delay_minutes"].mean()
    pct_delayed = (flights["departure_delay_minutes"] > 0).mean() * 100
    flights["ground_time_surplus"] = flights["scheduled_ground_time_minutes"] - flights["minimum_turn_minutes"]
    close_or_below = flights[flights["ground_time_surplus"] <= 5].shape[0]
    bags_agg = bags.groupby(["company_id","flight_number","scheduled_departure_date_local"]).agg(
        n_checked = ("bag_type", lambda s: (s=="Checked").sum()),
        n_transfer = ("bag_type", lambda s: (s=="Transfer").sum())
    ).reset_index()
    bags_agg["transfer_ratio"] = bags_agg["n_transfer"] / (bags_agg["n_checked"].replace(0, np.nan))
    avg_transfer_ratio = bags_agg["transfer_ratio"].mean()
    pnr_agg = pnr.groupby(["company_id","flight_number","scheduled_departure_date_local"]).agg(
        total_pax=("total_pax","sum"),
        lap_child_count=("lap_child_count","sum"),
        basic_economy_pax=("basic_economy_pax","sum"),
        is_stroller_user=("is_stroller_user","sum")
    ).reset_index()
    fl = flights.merge(pnr_agg, on=["company_id","flight_number","scheduled_departure_date_local"], how="left")
    fl["load_factor"] = fl["total_pax"] / fl["total_seats"]
    corr_load_delay = fl["load_factor"].corr(fl["departure_delay_minutes"])
    ssr_agg = pnr_remarks.groupby(["company_id","flight_number","scheduled_departure_date_local"]).agg(
        n_ssr=("special_service_request","count")
    ).reset_index()
    fl = fl.merge(ssr_agg, on=["company_id","flight_number","scheduled_departure_date_local"], how="left")
    from sklearn.linear_model import LinearRegression
    reg = LinearRegression()
    mask = fl[["departure_delay_minutes","load_factor","n_ssr"]].dropna().index
    if len(mask)>0:
        X = fl.loc[mask,["load_factor"]].fillna(0)
        y = fl.loc[mask,"departure_delay_minutes"].fillna(0)
        reg.fit(X, y)
        residuals = y - reg.predict(X)
        ssr_corr_with_resid = residuals.corr(fl.loc[mask,"n_ssr"])
    else:
        ssr_corr_with_resid = np.nan
    print("EDA Results:")
    print(f"Average departure delay (minutes): {avg_delay:.2f}")
    print(f"% flights departed later than scheduled: {pct_delayed:.1f}%")
    print(f"Flights with scheduled ground time <= minimum_turn + 5min: {close_or_below}")
    print(f"Average transfer/checked bag ratio (per flight, avg): {avg_transfer_ratio:.3f}")
    print(f"Correlation between load_factor and departure_delay: {corr_load_delay:.3f}")
    print(f"Correlation of SSR with delay residuals (controlling load): {ssr_corr_with_resid:.3f}")
    df = flights.merge(pnr_agg, on=["company_id","flight_number","scheduled_departure_date_local"], how="left")
    df = df.merge(bags_agg, on=["company_id","flight_number","scheduled_departure_date_local"], how="left")
    df = df.merge(ssr_agg, on=["company_id","flight_number","scheduled_departure_date_local"], how="left")
    df["total_pax"] = df["total_pax"].fillna(0)
    df["n_checked"] = df["n_checked"].fillna(0)
    df["n_transfer"] = df["n_transfer"].fillna(0)
    df["n_ssr"] = df["n_ssr"].fillna(0)
    df["total_seats"] = df["total_seats"].fillna(1)
    ensure_columns(df, {
        "scheduled_ground_time_minutes": np.nan,
        "minimum_turn_minutes": np.nan,
        "scheduled_departure_station_code": np.nan,
        "scheduled_arrival_station_code": np.nan,
        "fleet_type": "",
        "carrier": "",
        "scheduled_departure_datetime_local": pd.NaT,
        "scheduled_arrival_datetime_local": pd.NaT
    })
    df["load_factor"] = df["total_pax"] / df["total_seats"]
    df["transfer_share"] = df["n_transfer"] / (df["n_checked"] + df["n_transfer"] + 1e-6)
    df["hot_transfer_count"] = 0
    df["pax_per_bag"] = df["total_pax"] / (df["n_checked"].replace(0, np.nan))
    df["ground_time_surplus"] = df["scheduled_ground_time_minutes"] - df["minimum_turn_minutes"]
    df["short_turn_flag"] = (df["ground_time_surplus"] <= 5).astype(int)
    df["extreme_short_turn"] = (df["ground_time_surplus"] <= 0).astype(int)
    df["scheduled_departure_datetime_local"] = pd.to_datetime(df["scheduled_departure_datetime_local"], errors='coerce')
    df["scheduled_arrival_datetime_local"] = pd.to_datetime(df["scheduled_arrival_datetime_local"], errors='coerce')
    df["dep_hour"] = df["scheduled_departure_datetime_local"].dt.hour.fillna(-1).astype(int)
    df["dep_weekday"] = df["scheduled_departure_datetime_local"].dt.dayofweek.fillna(-1).astype(int)
    df["is_weekend"] = df["dep_weekday"].isin([5,6]).astype(int)
    df["is_peak_hour"] = df["dep_hour"].isin([6,7,8,9,16,17,18,19]).astype(int)
    df["basic_economy_share"] = df["basic_economy_pax"] / df["total_pax"].replace(0, np.nan)
    df["lap_child_share"] = df["lap_child_count"] / df["total_pax"].replace(0, np.nan)
    df["pax_per_seat"] = df["total_pax"] / df["total_seats"].replace(0, np.nan)
    df["ssr_per_pax"] = df["n_ssr"] / df["total_pax"].replace(0, np.nan)
    df["bags_per_pax"] = (df["n_checked"] + df["n_transfer"]) / df["total_pax"].replace(0, np.nan)
    df["transfer_intensity"] = df["n_transfer"] / (df["n_checked"] + 1)
    if "scheduled_departure_station_code" in df.columns:
        df["origin_daily_flights"] = df.groupby(["scheduled_departure_date_local","scheduled_departure_station_code"])['flight_number'].transform('count')
        df["origin_total_flights"] = df.groupby('scheduled_departure_station_code')['flight_number'].transform('count')
    else:
        df["origin_daily_flights"] = 0
        df["origin_total_flights"] = 0
    if "scheduled_arrival_station_code" in df.columns:
        df["dest_daily_flights"] = df.groupby(["scheduled_departure_date_local","scheduled_arrival_station_code"])['flight_number'].transform('count')
        df["dest_total_flights"] = df.groupby('scheduled_arrival_station_code')['flight_number'].transform('count')
    else:
        df["dest_daily_flights"] = 0
        df["dest_total_flights"] = 0
    try:
        ap = airports.rename(columns=lambda c: c.strip())
        possible_code_cols = [c for c in ap.columns if 'code' in c.lower() or 'icao' in c.lower() or 'iata' in c.lower()]
        if possible_code_cols:
            code_col = possible_code_cols[0]
            ap_small = ap[[code_col]].drop_duplicates().rename(columns={code_col: 'scheduled_arrival_station_code'})
            df = df.merge(ap_small, on='scheduled_arrival_station_code', how='left')
    except Exception:
        pass
    df["is_short_haul"] = ((df["scheduled_arrival_datetime_local"] - df["scheduled_departure_datetime_local"]).dt.total_seconds()/3600 < 2).astype(int)
    def compute_daily_score(group):
        g = group.copy()
        feature_dict = {}
        feature_dict["ground_pressure"] = -g["ground_time_surplus"].fillna(g["ground_time_surplus"].median())
        feature_dict["load_factor"] = g["load_factor"].fillna(0)
        feature_dict["transfer_share"] = g["transfer_share"].fillna(0)
        feature_dict["n_transfer"] = g["n_transfer"].fillna(0)
        feature_dict["n_checked"] = g["n_checked"].fillna(0)
        feature_dict["n_ssr"] = g["n_ssr"].fillna(0)
        feature_dict["basic_economy_pax"] = g["basic_economy_pax"].fillna(0)
        feature_dict["lap_child_count"] = g["lap_child_count"].fillna(0)
        feature_df = pd.DataFrame(feature_dict, index=g.index)
        def robust_z(x):
            med = x.median()
            mad = np.median(np.abs(x - med)) if len(x)>0 else 1.0
            mad = mad if mad>0 else (x.std() if x.std()>0 else 1.0)
            return (x - med) / mad
        feats_scaled = feature_df.apply(robust_z)
        weights = {
            "ground_pressure": 0.30,
            "load_factor": 0.20,
            "transfer_share": 0.15,
            "n_transfer": 0.05,
            "n_checked": 0.03,
            "n_ssr": 0.15,
            "basic_economy_pax": 0.06,
            "lap_child_count": 0.06
        }
        weighted = sum(feats_scaled[f] * weights.get(f,0) for f in feats_scaled.columns)
        weighted = weighted.fillna(weighted.median())
        minv, maxv = weighted.min(), weighted.max()
        if maxv - minv == 0:
            score_norm = pd.Series(0.5, index=weighted.index)
        else:
            score_norm = (weighted - minv) / (maxv - minv)
        g["difficulty_score_raw"] = weighted
        g["difficulty_score"] = score_norm
        g["difficulty_rank_within_day"] = g["difficulty_score"].rank(method="dense", ascending=False).astype(int)
        q = g["difficulty_score"].quantile([0.8, 0.3]).to_dict()
        def label(v):
            if v >= q[0.8]: return "Difficult"
            if v <= q[0.3]: return "Easy"
            return "Medium"
        g["difficulty_class"] = g["difficulty_score"].apply(label)
        for c in feats_scaled.columns:
            g[f"z_{c}"] = feats_scaled[c]
        return g
    out = df.groupby("scheduled_departure_date_local", group_keys=False).apply(compute_daily_score)
    deliver_cols = [
        "company_id","flight_number","scheduled_departure_date_local","scheduled_departure_datetime_local","scheduled_arrival_datetime_local",
        "scheduled_departure_station_code","scheduled_arrival_station_code","fleet_type","carrier",
        "total_seats","total_pax","load_factor",
        "scheduled_ground_time_minutes","minimum_turn_minutes","ground_time_surplus",
        "n_checked","n_transfer","transfer_share","n_ssr","basic_economy_pax","lap_child_count",
        "departure_delay_minutes",
        "difficulty_score","difficulty_rank_within_day","difficulty_class"
    ]
    deliver_cols = [c for c in deliver_cols if c in out.columns]
    out[deliver_cols].to_csv(OUTPUT_CSV, index=False)
    print(f"Written difficulty CSV to: {OUTPUT_CSV}")
    dest_diff = out[out["difficulty_class"]=="Difficult"].groupby("scheduled_arrival_station_code").agg(
        count_difficult=("flight_number","count"),
        avg_score=("difficulty_score","mean")
    ).sort_values("count_difficult", ascending=False).reset_index()
    print("\nTop destinations by difficult flight count (sample):")
    print(dest_diff.head(10))
    out.to_csv("flight_difficulty_full_debug.csv", index=False)
    PLOTS_DIR = "./plots"
    os.makedirs(PLOTS_DIR, exist_ok=True)
    def save_fig(fig, name):
        path = os.path.join(PLOTS_DIR, name)
        fig.savefig(path, bbox_inches='tight')
        plt.close(fig)
    if "difficulty_score" in out.columns:
        fig, ax = plt.subplots(figsize=(8,4))
        out["difficulty_score"].dropna().plot(kind='hist', bins=40, ax=ax, color='#2c7fb8')
        ax.set_title('Difficulty score distribution')
        ax.set_xlabel('difficulty_score')
        ax.set_ylabel('count')
        save_fig(fig, 'difficulty_score_hist.png')
    if not dest_diff.empty:
        topn = dest_diff.head(20)
        fig, ax = plt.subplots(figsize=(8,6))
        ax.barh(topn['scheduled_arrival_station_code'][::-1], topn['count_difficult'][::-1], color='#fdae6b')
        ax.set_xlabel('Difficult flight count')
        ax.set_title('Top destinations by difficult flight count')
        save_fig(fig, 'top_destinations_difficult_count.png')
    if 'load_factor' in fl.columns and 'departure_delay_minutes' in fl.columns:
        tmp = fl[['load_factor','departure_delay_minutes']].dropna()
        if not tmp.empty:
            x = tmp['load_factor'].values
            y = tmp['departure_delay_minutes'].values
            fig, ax = plt.subplots(figsize=(7,5))
            ax.scatter(x, y, alpha=0.4, s=10)
            try:
                coeffs = np.polyfit(x, y, deg=1)
                xs = np.linspace(x.min(), x.max(), 100)
                ax.plot(xs, np.polyval(coeffs, xs), color='red', linewidth=1)
            except Exception:
                pass
            ax.set_xlabel('load_factor')
            ax.set_ylabel('departure_delay_minutes')
            ax.set_title('Load factor vs departure delay')
            save_fig(fig, 'load_vs_delay_scatter.png')
    if 'scheduled_departure_date_local' in out.columns and 'difficulty_score' in out.columns:
        ts = out.groupby('scheduled_departure_date_local')['difficulty_score'].mean().reset_index()
        if not ts.empty:
            fig, ax = plt.subplots(figsize=(10,4))
            ax.plot(pd.to_datetime(ts['scheduled_departure_date_local']), ts['difficulty_score'], marker='o')
            ax.set_title('Mean daily difficulty score')
            ax.set_xlabel('date')
            ax.set_ylabel('mean difficulty_score')
            save_fig(fig, 'daily_mean_difficulty.png')
    try:
        out['is_difficult_top20'] = (out['difficulty_score'] >= out['difficulty_score'].quantile(0.8)).astype(int)
        feature_cols = [c for c in [
            'ground_time_surplus','load_factor','transfer_share','n_ssr','n_transfer','n_checked',
            'basic_economy_share','lap_child_share','ssr_per_pax','bags_per_pax','pax_per_seat','is_peak_hour',
            'short_turn_flag','extreme_short_turn','origin_daily_flights','dest_daily_flights'
        ] if c in out.columns]
        if feature_cols:
            from sklearn.ensemble import RandomForestClassifier
            rf = RandomForestClassifier(n_estimators=100, random_state=0)
            X = out[feature_cols].fillna(0)
            y = out['is_difficult_top20']
            rf.fit(X, y)
            importances = pd.Series(rf.feature_importances_, index=feature_cols).sort_values(ascending=False)
            fi_csv = os.path.join(PLOTS_DIR, 'feature_importances.csv')
            importances.to_csv(fi_csv, header=['importance'])
            fig, ax = plt.subplots(figsize=(8,6))
            sns.barplot(x=importances.values, y=importances.index, ax=ax, palette='viridis')
            ax.set_title('Feature importances (RF)')
            save_fig(fig, 'feature_importances.png')
    except Exception:
        pass
    try:
        from pptx import Presentation
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = 'Flight Difficulty Scoring & Insights'
        slide.placeholders[1].text = f'Author: {YOUR_NAME}\nGenerated'
        try:
            with open('PRESENTATION.md', 'r', encoding='utf8') as f:
                slides = f.read().split('\n\n')
            for s in slides:
                if not s.strip():
                    continue
                lines = [l.strip() for l in s.strip().splitlines() if l.strip()]
                title = lines[0][:60]
                body = lines[1:]
                slide_layout = prs.slide_layouts[1] if len(body)>0 else prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                try:
                    slide.shapes.title.text = title
                    if body:
                        tx = slide.shapes.placeholders[1].text_frame
                        tx.clear()
                        for b in body:
                            p = tx.add_paragraph()
                            p.text = b
                except Exception:
                    continue
        except FileNotFoundError:
            pass
        image_files = []
        try:
            for img in sorted(os.listdir(PLOTS_DIR)):
                if img.lower().endswith(('.png', '.jpg', '.jpeg')):
                    image_files.append(os.path.join(PLOTS_DIR, img))
        except Exception:
            image_files = []
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        for img_path in image_files:
            try:
                slide = prs.slides.add_slide(blank_layout)
                left = top = 0
                slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
            except Exception:
                continue
        prs.save('presentation.pptx')
    except Exception:
        pass
if __name__ == "__main__":
    try:
        main()
    except ModuleNotFoundError as e:
        print("A required Python package is not installed:", e.name)
        print("Install dependencies with: python -m pip install -r requirements.txt")
        raise
    except SystemExit as e:
        raise
    except Exception as e:
        print("An unexpected error occurred:", str(e))
        raise